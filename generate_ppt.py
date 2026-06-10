import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Color Palette
    DARK_TEAL = RGBColor(15, 118, 110)     # #0f766e
    LIGHT_TEAL = RGBColor(22, 163, 74)     # #16a34a
    WHITE = RGBColor(255, 255, 255)        # #ffffff
    LIGHT_BG = RGBColor(244, 247, 251)     # #f4f7fb
    CHARCOAL = RGBColor(23, 32, 42)        # #17202a
    MUTED_GREY = RGBColor(96, 112, 128)    # #607080
    ACCENT_MINT = RGBColor(183, 247, 230)   # #b7f7e6

    # Helper function to set slide background color
    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to add a clean title and subtitle/body text frame
    def add_title(slide, text, color, size=32, top=0.5):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color
        return top + 1.2

    # --- SLIDE 1: Title Slide (Dark Theme) ---
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1, DARK_TEAL)

    # Add decorative top banner block
    banner = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT_MINT
    banner.line.color.rgb = ACCENT_MINT

    # Main Title
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EXPENSE TRACKER"
    p.font.name = "Arial"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "A Full-Stack Expense Management Web Application"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_MINT
    p2.alignment = PP_ALIGN.LEFT

    # Project Info
    info_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(10), Inches(1.5))
    tf_info = info_box.text_frame
    p3 = tf_info.paragraphs[0]
    p3.text = "Tech Stack: HTML, CSS, JavaScript, Node.js, Express, MySQL"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE
    p3.font.italic = True

    p4 = tf_info.add_paragraph()
    p4.text = "Developer: Nikhil | Secure & Timezone-Independent Architecture"
    p4.font.name = "Arial"
    p4.font.size = Pt(14)
    p4.font.bold = True
    p4.font.color.rgb = ACCENT_MINT

    # --- SLIDE 2: Problem Statement (Light Theme) ---
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide2, LIGHT_BG)
    add_title(slide2, "The Problem Statement", DARK_TEAL)

    # Bullet points box
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True

    bullets = [
        "Traditional manual ledger records are tedious and prone to mathematical errors.",
        "Lack of real-time calculations (Total, Highest, Average) makes immediate budget analysis difficult.",
        "Server-client timezone differences often cause dates to shift unexpectedly (e.g. June 10th displaying as June 9th on different timezones).",
        "Form submissions without database input limits can trigger SQL failures or crash servers.",
        "Frontend applications without proper API validation crash during database downtime instead of recovering gracefully."
    ]

    for bullet in bullets:
        p = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
        p.text = "•  " + bullet
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = CHARCOAL
        p.space_after = Pt(20)

    # --- SLIDE 3: Proposed Solution (Light Theme) ---
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide3, LIGHT_BG)
    add_title(slide3, "The Proposed Solution", DARK_TEAL)

    content_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
    tf3 = content_box3.text_frame
    tf3.word_wrap = True

    solutions = [
        "A beautiful, responsive dashboard to log daily expenses with Title, Amount, Category, Date, and Optional Notes.",
        "Instant calculation of key metrics (Entries Count, Highest Expense, Average Expense, and Total Spent).",
        "Timezone-independent date tracking using raw strings (YYYY-MM-DD) bypasses browser date-shifting bugs.",
        "Strict client-side length constraints (maxlength 100 for titles, 255 for notes) matching MySQL column limits.",
        "Graceful error handling blocks front-end execution crashes during database disconnects and notifies the user."
    ]

    for sol in solutions:
        p = tf3.add_paragraph() if tf3.paragraphs[0].text else tf3.paragraphs[0]
        p.text = "✓  " + sol
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = CHARCOAL
        p.space_after = Pt(20)

    # --- SLIDE 4: Technology Stack (Light Theme) ---
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide4, LIGHT_BG)
    add_title(slide4, "System Architecture & Tech Stack", DARK_TEAL)

    # We will create 3 columns for Frontend, Backend, and Database
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(1.8)

    stacks = [
        ("Frontend UI", "• HTML5 & Semantic markup\n• Custom CSS3 (flexbox & grid)\n• Vanilla JS for DOM manipulation\n• Client-side input validation\n• Mobile-responsive layout"),
        ("Backend Server", "• Node.js runtime\n• Express.js framework\n• RESTful API architecture\n• dotenv for secure environment configs\n• CORS enabled for safe requests"),
        ("Database Layer", "• MySQL relational database\n• mysql2/promise for async pools\n• Structured SQL schemas\n• Decimal precision for amounts\n• Automatic timestamps")
    ]

    for idx, (title, content) in enumerate(stacks):
        left_pos = start_left + idx * (col_width + col_gap)
        
        # Draw a beautiful background card for each stack
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, col_width, Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = DARK_TEAL
        card.line.width = Pt(1.5)

        # Title of the stack card
        tb_title = slide4.shapes.add_textbox(left_pos, top_pos + Inches(0.2), col_width, Inches(0.8))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = DARK_TEAL
        p_title.alignment = PP_ALIGN.CENTER

        # Content of the stack card
        tb_content = slide4.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.9), col_width - Inches(0.4), Inches(3.5))
        tf_content = tb_content.text_frame
        tf_content.word_wrap = True
        
        for line in content.split("\n"):
            p_line = tf_content.add_paragraph() if tf_content.paragraphs[0].text else tf_content.paragraphs[0]
            p_line.text = line
            p_line.font.name = "Arial"
            p_line.font.size = Pt(14)
            p_line.font.color.rgb = CHARCOAL
            p_line.space_after = Pt(10)

    # --- SLIDE 5: Database Schema (Light Theme) ---
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide5, LIGHT_BG)
    add_title(slide5, "Database Schema Design", DARK_TEAL)

    # Display database script code and details
    desc_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    
    desc_points = [
        "Database name: expense_tracker",
        "A relational schema designed for optimal financial accuracy and indexing.",
        "Uses DECIMAL(10, 2) to prevent float-rounding errors on financial transactions.",
        "Uses MySQL AUTO_INCREMENT ID as primary key.",
        "Automatic created_at timestamp recording for transactions."
    ]
    for pt in desc_points:
        p = tf_desc.add_paragraph() if tf_desc.paragraphs[0].text else tf_desc.paragraphs[0]
        p.text = "• " + pt
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = CHARCOAL
        p.space_after = Pt(15)

    # Code block box
    code_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(1.8), Inches(6.3), Inches(4.8))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = CHARCOAL
    code_box.line.color.rgb = DARK_TEAL
    
    tb_code = slide5.shapes.add_textbox(Inches(6.3), Inches(1.9), Inches(6.1), Inches(4.6))
    tf_code = tb_code.text_frame
    tf_code.word_wrap = True
    sql_code = (
        "CREATE TABLE expenses (\n"
        "  id INT AUTO_INCREMENT PRIMARY KEY,\n"
        "  title VARCHAR(100) NOT NULL,\n"
        "  amount DECIMAL(10, 2) NOT NULL,\n"
        "  category VARCHAR(50) NOT NULL,\n"
        "  expense_date DATE NOT NULL,\n"
        "  note VARCHAR(255),\n"
        "  created_at TIMESTAMP DEFAULT\n"
        "               CURRENT_TIMESTAMP\n"
        ");"
    )
    p_code = tf_code.paragraphs[0]
    p_code.text = sql_code
    p_code.font.name = "Courier New"
    p_code.font.size = Pt(14)
    p_code.font.color.rgb = ACCENT_MINT

    # --- SLIDE 6: REST API Documentation (Light Theme) ---
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide6, LIGHT_BG)
    add_title(slide6, "RESTful API Endpoints", DARK_TEAL)

    # Add API Table
    rows = 4
    cols = 3
    left = Inches(0.8)
    top = Inches(2.0)
    width = Inches(11.7)
    height = Inches(4.0)
    
    table_shape = slide6.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set Column widths
    table.columns[0].width = Inches(2.2) # Method
    table.columns[1].width = Inches(3.5) # Route
    table.columns[2].width = Inches(6.0) # Description

    headers = ["HTTP Method", "Route / Endpoint", "Function / Description"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_TEAL
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    api_routes = [
        ("GET", "/api/expenses", "Fetches all transactions from the MySQL database, sorted in descending date order."),
        ("POST", "/api/expenses", "Adds a new expense. Validates all inputs (amount > 0, field lengths) before inserting."),
        ("DELETE", "/api/expenses/:id", "Deletes a single transaction matching the URL ID parameter.")
    ]

    for row_idx, (method, route, desc) in enumerate(api_routes):
        cell_m = table.cell(row_idx + 1, 0)
        cell_m.fill.solid()
        cell_m.fill.fore_color.rgb = WHITE
        p_m = cell_m.text_frame.paragraphs[0]
        p_m.text = method
        p_m.font.name = "Arial"
        p_m.font.size = Pt(15)
        p_m.font.bold = True
        p_m.font.color.rgb = LIGHT_TEAL
        p_m.alignment = PP_ALIGN.CENTER
        
        cell_r = table.cell(row_idx + 1, 1)
        cell_r.fill.solid()
        cell_r.fill.fore_color.rgb = WHITE
        p_r = cell_r.text_frame.paragraphs[0]
        p_r.text = route
        p_r.font.name = "Courier New"
        p_r.font.size = Pt(14)
        p_r.font.bold = True
        p_r.font.color.rgb = CHARCOAL
        
        cell_d = table.cell(row_idx + 1, 2)
        cell_d.fill.solid()
        cell_d.fill.fore_color.rgb = WHITE
        p_d = cell_d.text_frame.paragraphs[0]
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = CHARCOAL

    # --- SLIDE 7: Key Improvements & Debugging (Light Theme) ---
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide7, LIGHT_BG)
    add_title(slide7, "Engineering Highlights & Critical Fixes", DARK_TEAL)

    # Columns
    col_width7 = Inches(5.6)
    gap7 = Inches(0.533)
    start_left7 = Inches(0.8)
    top7 = Inches(1.8)

    fixes = [
        ("Timezone Safety Integration", 
         "• Enabled `dateStrings: true` in the connection pool config.\n"
         "• MySQL database returns date columns as raw strings (YYYY-MM-DD).\n"
         "• Prevents the server from auto-converting dates to UTC objects.\n"
         "• Eliminates the bug where dates displayed a day earlier or later for users in different timezones."),
        
        ("Robust Error Validation & Safety", 
         "• Form fields added `maxlength` bounds matching VARCHAR constraints.\n"
         "• Frontend checks `response.ok` before evaluating API data.\n"
         "• Catches server status codes (500, 404) to show warning dialogs.\n"
         "• Blocks the client script from crashing when database connections drop.")
    ]

    for idx, (title, content) in enumerate(fixes):
        left = start_left7 + idx * (col_width7 + gap7)
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top7, col_width7, Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = DARK_TEAL
        card.line.width = Pt(1.5)

        # Title
        tb_title = slide7.shapes.add_textbox(left, top7 + Inches(0.2), col_width7, Inches(0.8))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = DARK_TEAL

        # Content
        tb_content = slide7.shapes.add_textbox(left + Inches(0.2), top7 + Inches(0.9), col_width7 - Inches(0.4), Inches(3.6))
        tf_content = tb_content.text_frame
        tf_content.word_wrap = True
        
        for line in content.split("\n"):
            p_line = tf_content.add_paragraph() if tf_content.paragraphs[0].text else tf_content.paragraphs[0]
            p_line.text = line
            p_line.font.name = "Arial"
            p_line.font.size = Pt(14)
            p_line.font.color.rgb = CHARCOAL
            p_line.space_after = Pt(10)

    # --- SLIDE 8: Future Scope & Conclusion (Dark Theme) ---
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide8, DARK_TEAL)
    
    # Decorator top banner
    banner8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    banner8.fill.solid()
    banner8.fill.fore_color.rgb = ACCENT_MINT
    banner8.line.color.rgb = ACCENT_MINT

    add_title(slide8, "Future Project Enhancements", WHITE, size=40, top=1.0)

    content_box8 = slide8.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tf8 = content_box8.text_frame
    tf8.word_wrap = True

    future_scope = [
        "User Auth: Implementing JWT-based signup and login to support multiple user profiles.",
        "Data Visualizations: Integrating Chart.js or D3.js to show category-wise pie charts and month-on-month trend lines.",
        "Smart Budget Limits: Adding monthly budget caps with warnings when categories cross thresholds.",
        "Export Data: Support for downloading transaction history in Excel (.xlsx) and CSV format.",
        "Smart Recurrence: Automatically logging recurring transactions (Subscriptions, Rent, utilities) monthly."
    ]

    for pt in future_scope:
        p = tf8.add_paragraph() if tf8.paragraphs[0].text else tf8.paragraphs[0]
        p.text = "💡  " + pt
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(20)

    # Save presentation
    prs.save("Expense_Tracker_Presentation.pptx")
    print("Presentation created successfully: Expense_Tracker_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
